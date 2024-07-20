import os
import io
from flask import Flask, request, jsonify
from flask_cors import CORS
import google.generativeai as genai
import PyPDF2
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
import logging
from logging.handlers import RotatingFileHandler

load_dotenv()

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
    response.headers.add('Access-Control-Allow-Methods', 'GET,POST,OPTIONS')
    return response

app.config['GEMINI_API_KEY'] = os.getenv('GEMINI_API_KEY')

# Check if API key is loaded correctly
if not app.config['GEMINI_API_KEY']:
    raise ValueError("GEMINI_API_KEY is not set in the environment variables")

genai.configure(api_key=app.config['GEMINI_API_KEY'])
model = genai.GenerativeModel('gemini-1.5-flash')

def count_tokens(text):
    try:
        response = model.count_tokens(text)
        print(response.total_tokens)
        return response.total_tokens
    except Exception as e:
        app.logger.error(f"Error counting tokens: {e}")
        return 0

def summarize_text(text):
    try:
        prompt = f"Summarize the following text:\n\n{text}"
        if not text.strip():
            return ""
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        app.logger.error(f"Error summarizing text: {e}")
        return "Error summarizing text"

def extract_pdf_pages(file_stream):
    try:
        pdf_text = []
        reader = PyPDF2.PdfReader(file_stream)
        for page_num in range(len(reader.pages)):
            page = reader.pages[page_num]
            pdf_text.append(page.extract_text())
        return pdf_text
    except Exception as e:
        app.logger.error(f"Error extracting PDF pages: {e}")
        return []

def extract_ppt_pages(file_stream):
    try:
        ppt_text = []
        presentation = Presentation(file_stream)
        for slide in presentation.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + " "
            ppt_text.append(slide_text)
        return ppt_text
    except Exception as e:
        app.logger.error(f"Error extracting PPT pages: {e}")
        return []

def extract_docx_pages(file_stream):
    try:
        doc_text = []
        document = Document(file_stream)
        for para in document.paragraphs:
            doc_text.append(para.text)
        return doc_text
    except Exception as e:
        app.logger.error(f"Error extracting DOCX pages: {e}")
        return []

def process_document(file_stream, file_type, token_limit=20000):
    try:
        extractors = {
            'pdf': extract_pdf_pages,
            'ppt': extract_ppt_pages,
            'docx': extract_docx_pages
        }

        if file_type not in extractors:
            raise ValueError("Unsupported file type")

        extract_text = extractors[file_type]
        document_texts = extract_text(file_stream)
        combined_text = ""
        current_tokens = 0
        summary = ""

        for page_text in document_texts:
            if not page_text.strip():
                continue
            page_tokens = count_tokens(page_text)
            if page_tokens < 10:
                continue
            if current_tokens + page_tokens > token_limit:
                summary = summarize_text(combined_text)
                combined_text = page_text
                current_tokens = page_tokens
            else:
                combined_text += " " + page_text
                current_tokens += page_tokens

        if combined_text.strip():
            current_summary = summarize_text(combined_text)
            summary += current_summary

        print("Final Summary: ", summary)
        return summary
    except Exception as e:
        app.logger.error(f"Error processing document: {e}")
        return "Error processing document"

@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file part"}), 400
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No selected file"}), 400
        
        file_type = file.filename.rsplit('.', 1)[-1].lower()
        if file_type not in ['pdf', 'ppt', 'docx']:
            return jsonify({"error": "Unsupported file type"}), 400

        if file:
            summary = process_document(file, file_type)
            return jsonify({"summary": summary}), 200
    except Exception as e:
        app.logger.error(f"Error in upload_file endpoint: {e}")
        return jsonify({"error": "Internal Server Error"}), 500

@app.route('/chat', methods=['POST'])
def prompt():
    try:
        data = request.get_json()
        default = "Answer the question based on the context provided.\n\n"
        context = data['context']
        prompt = data['prompt']
        response = model.generate_content(default + "Question: \n\n" + prompt + "Context: \n\n" + context)
        print(response.text)
        return jsonify({"response": response.text})
    except Exception as e:
        app.logger.error(f"Error in chat endpoint: {e}")
        return jsonify({"error": "Internal Server Error"}), 500

if __name__ == '__main__':
    # Set up logging
    handler = RotatingFileHandler('app.log', maxBytes=10000, backupCount=1)
    handler.setLevel(logging.INFO)
    app.logger.addHandler(handler)
    
    app.run(host='0.0.0.0', port=8000, debug=True)
